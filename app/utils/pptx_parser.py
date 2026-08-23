import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_shape_data(shape, bullets, title_box, images, output_img_dir, rel_img_prefix, slide_num, img_counter):
    # 1. Text extraction
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            txt = p.text.strip()
            if txt:
                if not title_box['title']:
                    title_box['title'] = txt
                elif txt != title_box['title'] and txt not in bullets:
                    bullets.append(txt)

    # 2. Table extraction
    if shape.has_table:
        for row in shape.table.rows:
            row_txt = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_txt:
                line = " | ".join(row_txt)
                if line not in bullets:
                    bullets.append(line)

    # 3. Picture & Image extraction
    if output_img_dir and rel_img_prefix:
        try:
            if hasattr(shape, 'image') and shape.image:
                image = shape.image
                ext = image.ext or 'png'
                img_filename = f"slide_{slide_num}_img_{img_counter[0]}.{ext}"
                img_save_path = os.path.join(output_img_dir, img_filename)
                
                with open(img_save_path, 'wb') as f:
                    f.write(image.blob)
                
                rel_url = f"{rel_img_prefix}/{img_filename}"
                if rel_url not in images:
                    images.append(rel_url)
                    img_counter[0] += 1
        except Exception:
            pass

    # 4. Grouped shapes recursion
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            extract_shape_data(sub_shape, bullets, title_box, images, output_img_dir, rel_img_prefix, slide_num, img_counter)


def parse_pptx_slides(file_path, output_img_dir=None, rel_img_prefix=None):
    """
    Extracts slides, titles, bullet points, embedded images, and presenter notes from a .pptx file.
    Returns a list of dicts:
    [
        {
            "slide_number": 1,
            "title": "Slide Title",
            "bullets": ["Bullet point 1"],
            "images": ["/static/uploads/.../slide_1_img_1.png"],
            "notes": "Presenter notes"
        }
    ]
    """
    slides_data = []
    if not os.path.exists(file_path):
        return slides_data

    if output_img_dir:
        os.makedirs(output_img_dir, exist_ok=True)

    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides, start=1):
            title_box = {'title': None}
            bullets = []
            images = []
            notes = ""
            img_counter = [1]

            for shape in slide.shapes:
                if output_img_dir and rel_img_prefix:
                    extract_shape_data(shape, bullets, title_box, images, output_img_dir, rel_img_prefix, i, img_counter)
                else:
                    # Text only fallback
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            txt = p.text.strip()
                            if txt:
                                if not title_box['title']:
                                    title_box['title'] = txt
                                elif txt != title_box['title'] and txt not in bullets:
                                    bullets.append(txt)

            # Extract presenter notes if present
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()

            slides_data.append({
                "slide_number": i,
                "title": title_box['title'] or f"Slide {i}",
                "bullets": bullets,
                "images": images,
                "notes": notes
            })
    except Exception as e:
        print(f"Error parsing PPTX {file_path}: {e}")

    return slides_data
